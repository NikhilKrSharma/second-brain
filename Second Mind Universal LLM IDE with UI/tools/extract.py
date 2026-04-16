#!/usr/bin/env python3
"""
Extract documents into structured LLM-ready markdown.

Supports: PDF, DOCX, XLSX, PPTX

Usage:
    python tools/extract.py <path-to-file>
    python tools/extract.py raw/research/my-paper.pdf
    python tools/extract.py raw/work/report.docx
    python tools/extract.py raw/work/data.xlsx
    python tools/extract.py raw/work/slides.pptx

Dependencies (install only what you need):
    pip install pymupdf4llm          # for PDF
    pip install python-docx          # for DOCX
    pip install openpyxl             # for XLSX
    pip install python-pptx          # for PPTX

Output:
    <same-dir>/<stem>_extracted.md   — full structured markdown
    <same-dir>/<stem>_images/        — extracted images (if any, PDF/PPTX only)
"""

import sys
import re
import shutil
from pathlib import Path
from datetime import date
from typing import Optional


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}


# ---------------------------------------------------------------------------
# PDF extraction
# ---------------------------------------------------------------------------

def extract_pdf(file_path: Path, out_md: Path, image_dir: Path) -> dict:
    try:
        import pymupdf4llm
    except ImportError:
        print(
            "ERROR: pymupdf4llm is not installed.\n"
            "Run: pip install pymupdf4llm",
            file=sys.stderr,
        )
        sys.exit(1)

    md_text = pymupdf4llm.to_markdown(
        str(file_path),
        write_images=True,
        image_path=str(image_dir),
        image_format="png",
        image_size_limit=0.05,
        show_progress=False,
    )

    # Rewrite absolute image paths to relative
    abs_image_dir = str(image_dir.resolve())
    rel_image_dir = image_dir.name
    md_text = md_text.replace(abs_image_dir, rel_image_dir)

    # Get page count
    page_count = "unknown"
    try:
        import pymupdf
        doc = pymupdf.open(str(file_path))
        page_count = len(doc)
        doc.close()
    except Exception:
        pass

    image_count = len(list(image_dir.glob("*.png"))) if image_dir.exists() else 0

    header = (
        f"<!-- EXTRACTED FROM PDF -->\n"
        f"<!-- Source: {file_path.name} -->\n"
        f"<!-- Pages: {page_count} -->\n"
        f"<!-- Extracted: {date.today().isoformat()} -->\n\n"
    )

    out_md.write_text(header + md_text, encoding="utf-8")

    return {
        "words": len(re.findall(r"\w+", md_text)),
        "tables": md_text.count("|---|"),
        "images": image_count,
        "pages": page_count,
    }


# ---------------------------------------------------------------------------
# DOCX extraction
# ---------------------------------------------------------------------------

def extract_docx(file_path: Path, out_md: Path, image_dir: Path) -> dict:
    try:
        from docx import Document
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
    except ImportError:
        print(
            "ERROR: python-docx is not installed.\n"
            "Run: pip install python-docx",
            file=sys.stderr,
        )
        sys.exit(1)

    doc = Document(str(file_path))
    lines = []
    table_count = 0
    image_count = 0

    # Extract images
    image_dir.mkdir(parents=True, exist_ok=True)
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_count += 1
            img_name = f"image_{image_count}.png"
            img_path = image_dir / img_name
            try:
                img_path.write_bytes(rel.target_part.blob)
            except Exception:
                pass

    # Process document body
    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            # Paragraph
            para = None
            for p in doc.paragraphs:
                if p._element is element:
                    para = p
                    break
            if para is None:
                continue

            text = para.text.strip()
            if not text:
                lines.append("")
                continue

            style_name = para.style.name if para.style else ""

            if "Heading 1" in style_name:
                lines.append(f"# {text}")
            elif "Heading 2" in style_name:
                lines.append(f"## {text}")
            elif "Heading 3" in style_name:
                lines.append(f"### {text}")
            elif "Heading 4" in style_name:
                lines.append(f"#### {text}")
            elif "List" in style_name or "Bullet" in style_name:
                lines.append(f"- {text}")
            else:
                # Check for bold/italic runs
                parts = []
                for run in para.runs:
                    t = run.text
                    if not t:
                        continue
                    if run.bold and run.italic:
                        parts.append(f"***{t}***")
                    elif run.bold:
                        parts.append(f"**{t}**")
                    elif run.italic:
                        parts.append(f"*{t}*")
                    else:
                        parts.append(t)
                lines.append("".join(parts) if parts else text)

            lines.append("")

        elif tag == "tbl":
            # Table
            for table in doc.tables:
                if table._element is element:
                    table_count += 1
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        # Header row
                        lines.append("| " + " | ".join(rows[0]) + " |")
                        lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
                        for row in rows[1:]:
                            # Pad row if shorter than header
                            while len(row) < len(rows[0]):
                                row.append("")
                            lines.append("| " + " | ".join(row[:len(rows[0])]) + " |")
                        lines.append("")
                    break

    md_text = "\n".join(lines)

    header = (
        f"<!-- EXTRACTED FROM DOCX -->\n"
        f"<!-- Source: {file_path.name} -->\n"
        f"<!-- Extracted: {date.today().isoformat()} -->\n\n"
    )

    out_md.write_text(header + md_text, encoding="utf-8")

    return {
        "words": len(re.findall(r"\w+", md_text)),
        "tables": table_count,
        "images": image_count,
        "pages": "n/a",
    }


# ---------------------------------------------------------------------------
# XLSX extraction
# ---------------------------------------------------------------------------

def extract_xlsx(file_path: Path, out_md: Path, image_dir: Path) -> dict:
    try:
        from openpyxl import load_workbook
    except ImportError:
        print(
            "ERROR: openpyxl is not installed.\n"
            "Run: pip install openpyxl",
            file=sys.stderr,
        )
        sys.exit(1)

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    lines = []
    table_count = 0
    total_rows = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"## Sheet: {sheet_name}")
        lines.append("")

        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            # Skip completely empty rows
            if any(c.strip() for c in cells):
                rows.append(cells)

        if not rows:
            lines.append("*(empty sheet)*")
            lines.append("")
            continue

        total_rows += len(rows)
        table_count += 1

        # Normalize column count
        max_cols = max(len(r) for r in rows)
        for r in rows:
            while len(r) < max_cols:
                r.append("")

        # First row as header
        lines.append("| " + " | ".join(rows[0]) + " |")
        lines.append("| " + " | ".join(["---"] * max_cols) + " |")
        for row in rows[1:]:
            lines.append("| " + " | ".join(row[:max_cols]) + " |")
        lines.append("")

    sheet_count = len(wb.sheetnames)
    wb.close()

    md_text = "\n".join(lines)

    header = (
        f"<!-- EXTRACTED FROM XLSX -->\n"
        f"<!-- Source: {file_path.name} -->\n"
        f"<!-- Sheets: {sheet_count} -->\n"
        f"<!-- Extracted: {date.today().isoformat()} -->\n\n"
    )

    out_md.write_text(header + md_text, encoding="utf-8")

    return {
        "words": len(re.findall(r"\w+", md_text)),
        "tables": table_count,
        "images": 0,
        "pages": f"{sheet_count} sheets, {total_rows} rows",
    }


# ---------------------------------------------------------------------------
# PPTX extraction
# ---------------------------------------------------------------------------

def extract_pptx(file_path: Path, out_md: Path, image_dir: Path) -> dict:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print(
            "ERROR: python-pptx is not installed.\n"
            "Run: pip install python-pptx",
            file=sys.stderr,
        )
        sys.exit(1)

    prs = Presentation(str(file_path))
    lines = []
    image_count = 0
    table_count = 0

    image_dir.mkdir(parents=True, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {slide_num}")
        lines.append("")

        # Get slide title if available
        if slide.shapes.title and slide.shapes.title.text.strip():
            lines.append(f"### {slide.shapes.title.text.strip()}")
            lines.append("")

        for shape in slide.shapes:
            # Text content
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # Indent level → bullet depth
                        indent = para.level if para.level else 0
                        prefix = "  " * indent + "- " if indent > 0 else ""
                        lines.append(f"{prefix}{text}")
                lines.append("")

            # Tables
            if shape.has_table:
                table_count += 1
                tbl = shape.table
                rows = []
                for row in tbl.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    rows.append(cells)
                if rows:
                    max_cols = max(len(r) for r in rows)
                    for r in rows:
                        while len(r) < max_cols:
                            r.append("")
                    lines.append("| " + " | ".join(rows[0]) + " |")
                    lines.append("| " + " | ".join(["---"] * max_cols) + " |")
                    for row in rows[1:]:
                        lines.append("| " + " | ".join(row[:max_cols]) + " |")
                    lines.append("")

            # Images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                img_name = f"slide{slide_num}_img{image_count}.png"
                img_path = image_dir / img_name
                try:
                    img_path.write_bytes(shape.image.blob)
                    rel_image_dir = image_dir.name
                    lines.append(f"![{img_name}]({rel_image_dir}/{img_name})")
                    lines.append("")
                except Exception:
                    pass

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append(f"> **Speaker Notes:** {notes_text}")
                lines.append("")

        lines.append("---")
        lines.append("")

    md_text = "\n".join(lines)

    header = (
        f"<!-- EXTRACTED FROM PPTX -->\n"
        f"<!-- Source: {file_path.name} -->\n"
        f"<!-- Slides: {len(prs.slides)} -->\n"
        f"<!-- Extracted: {date.today().isoformat()} -->\n\n"
    )

    out_md.write_text(header + md_text, encoding="utf-8")

    return {
        "words": len(re.findall(r"\w+", md_text)),
        "tables": table_count,
        "images": image_count,
        "pages": f"{len(prs.slides)} slides",
    }


# ---------------------------------------------------------------------------
# Router
# ---------------------------------------------------------------------------

def extract(file_path: Path) -> tuple:
    if not file_path.exists():
        print(f"ERROR: File not found: {file_path}", file=sys.stderr)
        sys.exit(1)

    ext = file_path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        print(
            f"ERROR: Unsupported file type: {ext}\n"
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}",
            file=sys.stderr,
        )
        sys.exit(1)

    out_dir   = file_path.parent
    stem      = file_path.stem
    out_md    = out_dir / f"{stem}_extracted.md"
    image_dir = out_dir / f"{stem}_images"

    print(f"Extracting: {file_path}")

    if ext == ".pdf":
        stats = extract_pdf(file_path, out_md, image_dir)
    elif ext == ".docx":
        stats = extract_docx(file_path, out_md, image_dir)
    elif ext == ".xlsx":
        stats = extract_xlsx(file_path, out_md, image_dir)
    elif ext == ".pptx":
        stats = extract_pptx(file_path, out_md, image_dir)
    else:
        print(f"ERROR: No extractor for {ext}", file=sys.stderr)
        sys.exit(1)

    print(f"Done.")
    print(f"  Output:  {out_md}")
    print(f"  Words:   {stats['words']:,}")
    print(f"  Tables:  {stats['tables']}")
    print(f"  Images:  {stats['images']}")
    if stats.get("pages"):
        print(f"  Pages:   {stats['pages']}")

    return out_md, image_dir


def copy_images_to_wiki(image_dir: Path, slug: str) -> Optional[Path]:
    """Copy extracted images into wiki/assets/images/<slug>/ and return the dest dir."""
    if not image_dir.exists() or not any(image_dir.iterdir()):
        return None

    repo_root = Path(__file__).parent.parent
    dest = repo_root / "wiki" / "assets" / "images" / slug
    dest.mkdir(parents=True, exist_ok=True)

    copied = 0
    for img in image_dir.iterdir():
        if img.is_file():
            shutil.copy2(img, dest / img.name)
            copied += 1

    print(f"  Copied {copied} image(s) → {dest.relative_to(repo_root)}")
    return dest


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Extract documents into structured LLM-ready markdown.",
        epilog=f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}",
    )
    parser.add_argument("file", help="Path to the file to extract")
    parser.add_argument(
        "--copy-images",
        metavar="SLUG",
        help="After extraction, copy images into wiki/assets/images/<SLUG>/",
    )
    args = parser.parse_args()

    out_md, image_dir = extract(Path(args.file))

    if args.copy_images:
        copy_images_to_wiki(image_dir, args.copy_images)
